"""
content_extractor.py — Module d'extraction de contenu textuel depuis différents types de fichiers.

Stratégie par type :
    PDF  → pymupdf4llm (extraction Markdown structurée + OCR automatique ou forcé)
    XLSX → openpyxl (lecture feuille par feuille, valeurs brutes)
    PPTX → python-pptx (extraction des textes de slides et tableaux)
    DOCX → python-docx (paragraphes + tableaux)
    CSV  → lecture directe avec détection d'encodage
    MP4  → ffmpeg (extraction audio) + faster-whisper (transcription STT)
    Audio → faster-whisper directement

Utilisation :
    from lib.core.content_extractor import extract_content, get_supported_extensions

    text = extract_content("/tmp/doc.pdf")           # extension déduite du chemin
    text = extract_content("/tmp/upload", ".pdf")    # extension fournie explicitement
"""

import os
import subprocess
import tempfile
import platform
import re
from pathlib import Path


# ──────────────────────────────────────────────────────────────
# CONFIGURATION
# ──────────────────────────────────────────────────────────────

# Longueur minimale de texte utile après nettoyage des artefacts Markdown
MIN_USEFUL_TEXT_LENGTH = 50

# Correspondance extension → extracteur interne
SUPPORTED_EXTENSIONS: dict[str, str] = {
    ".pdf":  "pdf",
    ".xlsx": "xlsx",
    ".xls":  "xlsx",
    ".pptx": "pptx",
    ".docx": "docx",
    ".csv":  "csv",
    ".mp4":  "mp4",
    ".mp3":  "audio",
    ".wav":  "audio",
    ".m4a":  "audio",
    ".webm": "audio",
    ".ogg":  "audio",
    # Formats image — OCR via pymupdf
    ".png":  "image",
    ".jpg":  "image",
    ".jpeg": "image",
    ".bmp":  "image",
    ".tiff": "image",
    ".tif":  "image",
    ".webp": "image",
    ".gif":  "image",
}


def _safe_path(path: str) -> str:
    """Retourne un chemin absolu compatible Windows Long Path (préfixe \\\\?\\)."""
    abs_path = os.path.abspath(path)
    if platform.system() == "Windows" and not abs_path.startswith("\\\\?\\"):
        abs_path = "\\\\?\\" + abs_path
    return abs_path


# ──────────────────────────────────────────────────────────────
# EXTRACTEURS PAR TYPE
# ──────────────────────────────────────────────────────────────

def _is_meaningful_text(text: str) -> bool:
    """
    Vérifie qu'un texte extrait contient du contenu exploitable.

    Les PDF scannés sans couche OCR renvoient souvent uniquement
    des séquences de tirets ou de syntaxe Markdown vide via pymupdf4llm.
    On nettoie ces artefacts avant de mesurer la longueur utile.
    """
    cleaned = re.sub(r"-{2,}", "", text)         # séquences de tirets
    cleaned = re.sub(r"\s+", "", cleaned)         # espaces / sauts de ligne
    cleaned = re.sub(r"[#*_|>~`]", "", cleaned)  # syntaxe Markdown résiduelle
    return len(cleaned) >= MIN_USEFUL_TEXT_LENGTH


def extract_pdf(file_path: str) -> str | None:
    """
    Extrait le contenu d'un PDF en Markdown structuré via pymupdf4llm.

    Stratégie en deux passes :
        1. Extraction standard avec auto-détection OCR (rapide pour PDF textuels).
        2. Si le résultat est inexploitable → forcer l'OCR complet (PDF scannés).

    Returns:
        Texte Markdown ou None si aucun contenu exploitable n'est trouvé.
    """
    try:
        import pymupdf4llm

        safe = _safe_path(file_path)

        # ── Passe 1 : extraction avec OCR auto ──
        md_text: str = pymupdf4llm.to_markdown(safe, ocr_language="fra")
        if md_text and md_text.strip() and _is_meaningful_text(md_text):
            return md_text

        # ── Passe 2 : OCR forcé (couche texte absente ou corrompue) ──
        print(f"  [PDF] Passe 1 non exploitable — OCR forcé : {os.path.basename(file_path)}")
        md_text = pymupdf4llm.to_markdown(safe, force_ocr=True, ocr_language="fra")

        if not md_text or not md_text.strip():
            print(f"  [PDF] Aucun texte extrait : {file_path}")
            return None

        if not _is_meaningful_text(md_text):
            print(f"  [PDF] Contenu non exploitable même après OCR forcé : {file_path}")
            return None

        return md_text

    except ImportError:
        print("  [PDF] pymupdf4llm non installé. Exécuter : pip install pymupdf4llm")
        return None
    except Exception as exc:
        print(f"  [PDF] Erreur extraction : {exc}")
        return None


def extract_xlsx(file_path: str) -> str | None:
    """
    Extrait le contenu textuel d'un fichier Excel (.xlsx / .xls).
    Parcourt toutes les feuilles cellule par cellule.
    """
    try:
        from openpyxl import load_workbook

        wb = load_workbook(_safe_path(file_path), read_only=True, data_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Feuille : {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    parts.append(row_text)

        wb.close()
        return "\n".join(parts) if parts else None

    except ImportError:
        print("  [XLSX] openpyxl non installé. Exécuter : pip install openpyxl")
        return None
    except Exception as exc:
        print(f"  [XLSX] Erreur extraction : {exc}")
        return None


def extract_csv(file_path: str) -> str | None:
    """
    Lit un fichier CSV et retourne son contenu brut.
    Tente plusieurs encodages courants (utf-8, latin-1, cp1252).
    """
    try:
        for encoding in ("utf-8", "latin-1", "cp1252"):
            try:
                with open(_safe_path(file_path), "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        print(f"  [CSV] Aucun encodage compatible trouvé pour : {file_path}")
        return None
    except Exception as exc:
        print(f"  [CSV] Erreur extraction : {exc}")
        return None


def extract_pptx(file_path: str) -> str | None:
    """
    Extrait le texte de toutes les diapositives d'un fichier PowerPoint (.pptx).
    Inclut les zones de texte libres et les tableaux.
    """
    try:
        from pptx import Presentation

        prs = Presentation(_safe_path(file_path))
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []

            for shape in slide.shapes:
                # Zones de texte libre
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                # Tableaux intégrés
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            slide_texts.append(row_text)

            if slide_texts:
                parts.append(f"--- Slide {slide_num} ---")
                parts.extend(slide_texts)

        return "\n".join(parts) if parts else None

    except ImportError:
        print("  [PPTX] python-pptx non installé. Exécuter : pip install python-pptx")
        return None
    except Exception as exc:
        print(f"  [PPTX] Erreur extraction : {exc}")
        return None


def extract_docx(file_path: str) -> str | None:
    """
    Extrait le texte d'un fichier Word (.docx).
    Inclut paragraphes narratifs et cellules de tableaux.
    """
    try:
        from docx import Document

        doc = Document(_safe_path(file_path))
        parts: list[str] = []

        # Paragraphes
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)

        # Tableaux
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)

        return "\n".join(parts) if parts else None

    except ImportError:
        print("  [DOCX] python-docx non installé. Exécuter : pip install python-docx")
        return None
    except Exception as exc:
        print(f"  [DOCX] Erreur extraction : {exc}")
        return None


def extract_video_audio(file_path: str) -> str | None:
    """
    Transcrit un fichier audio ou vidéo en texte via faster-whisper.

    Flux :
        - Fichiers vidéo (.mp4, .webm, .ogg) → extraction audio WAV avec ffmpeg
        - Fichiers audio purs (.mp3, .wav, .m4a) → transcription directe
    """
    try:
        from faster_whisper import WhisperModel

        ext = Path(file_path).suffix.lower()
        audio_only_extensions = {".mp3", ".wav", ".m4a"}

        if ext in audio_only_extensions:
            audio_path = _safe_path(file_path)
            tmp_dir = None
        else:
            # Extraction audio via ffmpeg
            tmp_dir = tempfile.mkdtemp()
            audio_path = os.path.join(tmp_dir, "audio.wav")
            cmd = [
                "ffmpeg", "-i", _safe_path(file_path),
                "-vn",               # pas de piste vidéo
                "-acodec", "pcm_s16le",
                "-ar", "16000",      # fréquence 16 kHz requise par Whisper
                "-ac", "1",          # mono
                "-y",                # écraser sans confirmation
                audio_path,
            ]
            result = subprocess.run(cmd, capture_output=True, timeout=600)
            if result.returncode != 0:
                print(f"  [FFMPEG] Erreur extraction audio : {result.stderr.decode()[:500]}")
                return None

        # Transcription Whisper (modèle "small" : compromis vitesse / qualité)
        model = WhisperModel("small", device="cpu", compute_type="int8")
        segments, _info = model.transcribe(audio_path, language="fr")
        transcript = "\n".join(segment.text.strip() for segment in segments)

        # Nettoyage du répertoire temporaire
        if tmp_dir:
            try:
                os.remove(audio_path)
                os.rmdir(tmp_dir)
            except OSError:
                pass

        return transcript if transcript.strip() else None

    except ImportError:
        print("  [AUDIO] faster-whisper non installé. Exécuter : pip install faster-whisper")
        return None
    except Exception as exc:
        print(f"  [AUDIO] Erreur transcription : {exc}")
        return None


def extract_image(file_path: str) -> str | None:
    """
    Extrait le texte d'une image (PNG, JPG, BMP, TIFF, WEBP) via OCR pymupdf.

    PyMuPDF peut ouvrir directement les fichiers image comme des documents
    mono-page et y appliquer l'OCR Tesseract (langage français).
    Ce résultat sert de fallback textuel ; pour une analyse approfondie
    d'un plan BTP, préférer l'analyse vision via Gemini (smart_btp_agent).

    Returns:
        Texte OCR extrait, ou None si aucun résultat exploitable.
    """
    try:
        import pymupdf  # fitz ≥ 1.23

        safe = _safe_path(file_path)
        doc = pymupdf.open(safe)
        parts: list[str] = []

        for page in doc:
            # Passe 1 : texte natif (images avec couche texte embarquée)
            text = page.get_text("text").strip()
            if text and _is_meaningful_text(text):
                parts.append(text)
                continue

            # Passe 2 : OCR Tesseract (images scannées ou photographiées)
            try:
                tp = page.get_textpage_ocr(flags=0, language="fra", dpi=300, full=True)
                text = page.get_text("text", textpage=tp).strip()
                if text and _is_meaningful_text(text):
                    parts.append(text)
            except Exception as ocr_exc:
                print(f"  [IMAGE] OCR page {page.number} échoué : {ocr_exc}")

        doc.close()
        return "\n".join(parts) if parts else None

    except ImportError:
        print("  [IMAGE] pymupdf non installé. Exécuter : pip install pymupdf")
        return None
    except Exception as exc:
        print(f"  [IMAGE] Erreur extraction : {exc}")
        return None


# ──────────────────────────────────────────────────────────────
# DISPATCHER PRINCIPAL
# ──────────────────────────────────────────────────────────────

# Table de routage vers les extracteurs
_EXTRACTORS: dict[str, callable] = {
    "pdf":   extract_pdf,
    "xlsx":  extract_xlsx,
    "csv":   extract_csv,
    "pptx":  extract_pptx,
    "docx":  extract_docx,
    "mp4":   extract_video_audio,
    "audio": extract_video_audio,
    "image": extract_image,
}


def extract_content(file_path: str, extension: str | None = None) -> str | None:
    """
    Point d'entrée principal : extrait le contenu textuel d'un fichier
    selon son extension.

    Args:
        file_path : Chemin local vers le fichier à traiter.
        extension : Extension explicite (ex: '.pdf'). Si None, déduite du chemin.

    Returns:
        Texte extrait (str) ou None si l'extraction échoue ou si le type
        n'est pas supporté.
    """
    # ── Normalisation de l'extension ──
    if extension is None:
        ext = Path(file_path).suffix.lower()
    else:
        ext = extension.lower()
        if not ext.startswith("."):
            ext = f".{ext}"

    # ── Vérification de l'existence du fichier ──
    if not os.path.exists(_safe_path(file_path)):
        print(f"  [EXTRACT] Fichier introuvable : {file_path}")
        return None

    # ── Résolution de l'extracteur ──
    ext_type = SUPPORTED_EXTENSIONS.get(ext)
    if ext_type is None:
        print(f"  [EXTRACT] Extension non supportée : {ext}")
        return None

    extractor = _EXTRACTORS.get(ext_type)
    if extractor is None:
        return None

    print(f"  [EXTRACT] {ext.upper()} → {ext_type} : {os.path.basename(file_path)}")
    return extractor(file_path)


def get_supported_extensions() -> list[str]:
    """Retourne la liste des extensions de fichiers prises en charge."""
    return list(SUPPORTED_EXTENSIONS.keys())
